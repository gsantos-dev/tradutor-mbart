import os
from pptx import Presentation


def salvar_texto_unificado(caminho_entrada, caminho_saida, nome_saida):
    if not os.path.exists(caminho_saida):
        os.makedirs(caminho_saida)

    texto_total = ""

    for arquivo in sorted(os.listdir(caminho_entrada)):
        if arquivo.endswith(".pptx"):
            prs = Presentation(os.path.join(caminho_entrada, arquivo))
            texto_ppt = f"\n### {arquivo} ###\n"

            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texto_ppt += shape.text + "\n"

                    elif shape.has_table:
                        for row in shape.table.rows:
                            texto_ppt += "\t".join([cell.text.strip() for cell in row.cells]) + "\n"

            texto_total += texto_ppt + "\n"
    caminho_txt = os.path.join(caminho_saida, nome_saida)

    with open(caminho_txt, "w", encoding="utf-8") as f:
        f.write(texto_total.strip())

    print(f"Texto consolidado salvo: {caminho_txt}")


caminho_ppt_pt = "/PPTs_pt"
caminho_ppt_en = "/PPTs_en"
caminho_saida_pt = "/Texto_pt"
caminho_saida_en = "/Texto_en"


salvar_texto_unificado(caminho_ppt_pt, caminho_saida_pt, "Texto_Completo_PT.txt")
salvar_texto_unificado(caminho_ppt_en, caminho_saida_en, "Texto_Completo_EN.txt")


