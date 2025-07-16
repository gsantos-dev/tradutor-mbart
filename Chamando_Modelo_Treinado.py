import os
import torch
from transformers import MBartForConditionalGeneration, MBart50TokenizerFast
from pptx import Presentation

def carregar_modelo(model_path):
    tokenizer = MBart50TokenizerFast.from_pretrained(model_path)
    model = MBartForConditionalGeneration.from_pretrained(model_path)
    tokenizer.src_lang = "pt_XX"
    tokenizer.tgt_lang = "en_XX"
    model.to("cpu")
    return model, tokenizer

def traduzir_texto(texto, model, tokenizer):
    inputs = tokenizer(texto, return_tensors="pt", padding=True, truncation=True, max_length=512)
    translated_tokens = model.generate(**inputs)
    return tokenizer.batch_decode(translated_tokens, skip_special_tokens=True)[0]

def traduzir_ppt(input_ppt, output_ppt, model, tokenizer):
    prs = Presentation(input_ppt)

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                shape.text = traduzir_texto(shape.text, model, tokenizer)
            elif shape.has_table:
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            cell.text = traduzir_texto(cell.text, model, tokenizer)

    prs.save(output_ppt)
    print(f"✅ Tradução concluída: {output_ppt}")

def processar_pasta(input_folder, output_folder, model, tokenizer):
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    for file_name in os.listdir(input_folder):
        if file_name.endswith(".pptx"):
            input_path = os.path.join(input_folder, file_name)
            output_path = os.path.join(output_folder, file_name.replace(".pptx", "_traduzido.pptx"))

            print(f"🔄 Traduzindo: {file_name}...")
            traduzir_ppt(input_path, output_path, model, tokenizer)


model_path = "C:/Users/Guilherme/Desktop/Projetos/Projeto_Tradutor/Modelo_Mbart_Treinado"
input_folder = "C:/Users/Guilherme/Desktop/Projetos/Projeto_Tradutor/Input_PPTS"
output_folder = "C:/Users/Guilherme/Desktop/Projetos/Projeto_Tradutor/Outputs_PPTS"


model, tokenizer = carregar_modelo(model_path)
processar_pasta(input_folder, output_folder, model, tokenizer)
